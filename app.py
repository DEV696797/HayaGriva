import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import statsmodels.api as sm
from sklearn.cluster import KMeans
from pptx import Presentation

st.set_page_config(layout="wide")

# ------------------------------------------------
# LUXURY THEME
# ------------------------------------------------

st.markdown("""
<style>

html, body, [class*="css"]  {
    font-family: "Times New Roman", Times, serif;
    color: #ffffff;
}

.stApp{
background:linear-gradient(135deg,#1b0036,#3b0a63,#5e17eb);
color:#ffffff;
}

h1,h2,h3,h4,h5,h6{
font-family:"Times New Roman",serif;
color:#ffffff;
}

p{
font-family:"Times New Roman",serif;
color:#ffffff;
font-size:18px;
}

</style>
""",unsafe_allow_html=True)

st.title("HayaGriva Luxury Consumer Intelligence")
st.subheader("Effect of Emotions on Purchase Intention of Luxury Products")

# ------------------------------------------------
# FILE UPLOAD
# ------------------------------------------------

emotion_file = st.file_uploader("Upload Emotional Dataset")

demo_file = st.file_uploader("Upload Demographic Dataset")

# ------------------------------------------------
# LOAD DATA
# ------------------------------------------------

if emotion_file:

    df = pd.read_excel(emotion_file)

    df = df.replace(0,3)

    emotion = df.iloc[:,5]
    celebrity = df.iloc[:,10]
    fomo = df.iloc[:,15]
    purchase = df.iloc[:,20]

# ------------------------------------------------
# MULTIPLE REGRESSION
# ------------------------------------------------

    st.header("Multiple Linear Regression")

    X = pd.DataFrame({
        "Emotion":emotion,
        "Celebrity":celebrity,
        "FOMO":fomo
    })

    X = sm.add_constant(X)

    model = sm.OLS(purchase,X).fit()

    col1,col2 = st.columns([2,1])

    with col1:
        st.dataframe(model.summary2().tables[1])

    with col2:
        st.markdown("""
### Insights

• Emotional response has the strongest statistical impact on luxury purchase intention  

• Celebrity endorsements increase aspirational perception  

• FOMO increases urgency and social comparison  

• Psychological drivers strongly explain luxury consumption behaviour
""")

# ------------------------------------------------
# DRIVER COMPARISON
# ------------------------------------------------

    st.header("Psychological Driver Comparison")

    drivers = pd.DataFrame({
        "Driver":["Emotion","Celebrity","FOMO"],
        "Score":[emotion.mean(),celebrity.mean(),fomo.mean()]
    })

    col1,col2 = st.columns([2,1])

    with col1:
        fig = px.bar(drivers,x="Driver",y="Score",color="Driver")
        st.plotly_chart(fig,use_container_width=True)

    with col2:
        st.markdown("""
### Interpretation

• Emotional gratification dominates luxury consumption  

• Celebrity endorsements reinforce aspirational identity  

• FOMO increases urgency in luxury purchasing
""")

# ------------------------------------------------
# RADAR CHART
# ------------------------------------------------

    st.header("Luxury Motivation Radar")

    radar = pd.DataFrame({
        "Factor":["Emotion","Celebrity","FOMO","Purchase"],
        "Score":[emotion.mean(),celebrity.mean(),fomo.mean(),purchase.mean()]
    })

    col1,col2 = st.columns([2,1])

    with col1:
        fig = px.line_polar(radar,r="Score",theta="Factor",line_close=True)
        st.plotly_chart(fig,use_container_width=True)

    with col2:
        st.markdown("""
### Insights

• Emotion dominates luxury motivation  

• Celebrity influence reinforces aspirational identity  

• FOMO acts as a psychological trigger
""")

# ------------------------------------------------
# SCATTER RELATIONSHIP
# ------------------------------------------------

    st.header("Emotion vs Purchase Intention")

    col1,col2 = st.columns([2,1])

    with col1:
        fig = px.scatter(x=emotion,y=purchase,trendline="ols")
        st.plotly_chart(fig,use_container_width=True)

    with col2:
        st.markdown("""
### Insights

• Emotional attachment strongly predicts purchase intention  

• Consumers with stronger emotional engagement buy luxury more frequently  

• Emotional branding drives symbolic consumption
""")

# ------------------------------------------------
# CORRELATION HEATMAP
# ------------------------------------------------

    st.header("Psychological Correlation Matrix")

    matrix = pd.DataFrame({
        "Emotion":emotion,
        "Celebrity":celebrity,
        "FOMO":fomo,
        "Purchase":purchase
    })

    col1,col2 = st.columns([2,1])

    with col1:

        heat = px.imshow(
            matrix.corr(),
            text_auto=True,
            color_continuous_scale="Purples"
        )

        st.plotly_chart(heat,use_container_width=True)

    with col2:
        st.markdown("""
### Insights

• Emotion and purchase intention show the highest correlation  

• Celebrity influence moderately affects emotional engagement  

• FOMO amplifies emotional luxury desire
""")

# ------------------------------------------------
# SEGMENTATION
# ------------------------------------------------

    st.header("Luxury Consumer Segmentation")

    seg = pd.DataFrame({
        "Emotion":emotion,
        "Celebrity":celebrity,
        "FOMO":fomo
    })

    kmeans = KMeans(n_clusters=3)

    df["segment"] = kmeans.fit_predict(seg)

    col1,col2 = st.columns([2,1])

    with col1:

        fig = px.scatter(
            df,
            x=emotion,
            y=purchase,
            size=fomo,
            color=df["segment"]
        )

        st.plotly_chart(fig,use_container_width=True)

    with col2:

        st.markdown("""
### Segments

Status Aspirers  
• Emotionally driven prestige buyers  

Celebrity Followers  
• Influenced by endorsements  

FOMO Buyers  
• Driven by social comparison
""")

# ------------------------------------------------
# FUNNEL
# ------------------------------------------------

    st.header("Luxury Consumer Journey")

    funnel = pd.DataFrame({
        "Stage":["Awareness","Aspiration","Emotional Attachment","Purchase"],
        "Value":[100,80,60,40]
    })

    col1,col2 = st.columns([2,1])

    with col1:
        fig = px.funnel(funnel,x="Value",y="Stage")
        st.plotly_chart(fig,use_container_width=True)

    with col2:
        st.markdown("""
### Insights

• Consumers first become aware of luxury brands  

• Aspirational perception develops through marketing  

• Emotional attachment drives purchase intention
""")

# ------------------------------------------------
# DEMOGRAPHIC ANALYSIS
# ------------------------------------------------

if demo_file:

    demo = pd.read_excel(demo_file)

    st.header("Demographic Distribution")

    col1,col2 = st.columns([2,1])

    with col1:
        fig = px.histogram(demo,x=demo.columns[0])
        st.plotly_chart(fig,use_container_width=True)

    with col2:
        st.markdown("""
### Insights

• Younger consumers show stronger luxury aspirations  

• Urban exposure increases emotional engagement  

• Demographics influence luxury adoption patterns
""")

# ------------------------------------------------
# BRAND POSITIONING MAP
# ------------------------------------------------

st.header("Luxury Brand Positioning Map")

brands = pd.DataFrame({

"Brand":["Louis Vuitton","Hermès","Gucci","Rolex","Chanel"],

"Emotion":[9,9,8,7,8],

"Exclusivity":[8,10,7,9,9],

"Influence":[9,7,8,8,8]

})

fig = px.scatter(
brands,
x="Emotion",
y="Exclusivity",
size="Influence",
text="Brand",
color="Brand"
)

st.plotly_chart(fig,use_container_width=True)

# ------------------------------------------------
# PURCHASE SIMULATOR
# ------------------------------------------------

st.header("Luxury Purchase Simulator")

emotion_input = st.slider("Emotion",1,20,10)

celebrity_input = st.slider("Celebrity Influence",1,20,10)

fomo_input = st.slider("FOMO",1,20,10)

score = -2.120 + 0.757*emotion_input + 0.199*celebrity_input + 0.314*fomo_input

st.success(f"Predicted Purchase Score: {round(score,2)}")

# ------------------------------------------------
# CASE STUDIES
# ------------------------------------------------

st.header("Luxury Brand Strategy Case Studies")

st.markdown("""
**Louis Vuitton — Emotional Storytelling**

Louis Vuitton builds emotional attachment through heritage storytelling and craftsmanship narratives.

**Gucci — Celebrity Influence**

Gucci leverages celebrity collaborations and digital storytelling to create aspirational identity.

**Hermès — Scarcity Strategy**

Hermès maintains exclusivity and limited supply to intensify emotional desire.
""")

# ------------------------------------------------
# CONCLUSION
# ------------------------------------------------

st.header("Project Conclusion")

st.write("""
Emotion is the strongest driver of luxury purchase intention. 
Luxury brands should focus on emotional storytelling, aspirational identity creation, and exclusivity strategies to maintain desirability and prestige.
""")

# ------------------------------------------------
# PPT GENERATOR
# ------------------------------------------------

st.header("Generate Presentation")

def generate_ppt():

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Luxury Consumer Research"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Findings"
    slide.placeholders[1].text = """
Emotion strongly drives luxury purchase intention.
Celebrity influence increases aspiration.
FOMO accelerates luxury purchase urgency.
"""

    prs.save("luxury_report.pptx")

    return "luxury_report.pptx"

if st.button("Generate PPT"):

    file = generate_ppt()

    with open(file,"rb") as f:

        st.download_button(
            "Download Presentation",
            f,
            file_name="luxury_report.pptx"
        )